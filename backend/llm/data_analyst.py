"""
backend/llm/data_analyst.py

データ読込後のLLM対話型分析プラン立案サービス。
GGUFProvider (Bonsai 8B) を使用し、多輪対話で分析方針を提案する。
記述子推奨機能付き。
"""

from __future__ import annotations

import json
import logging
from typing import Any

import pandas as pd

from backend.llm.descriptor_knowledge import (
    PROPERTY_CATEGORIES,
    build_descriptor_recommendation_prompt,
    find_matching_properties,
)

logger = logging.getLogger(__name__)


def read_document_content(filename: str, content: bytes | None = None) -> dict:
    """
    ドキュメントファイルからテキストとテーブルを抽出する。

    Args:
        filename: ファイル名（拡張子判定用）
        content: ファイルのバイト列（アップロード時）。Noneの場合はファイルを読み込む

    Returns:
        {
            "text": str,           # 抽出されたテキスト
            "tables": list[pd.DataFrame],  # 抽出されたテーブル
            "metadata": dict        # メタデータ（ページ数、スライド数など）
        }
    """
    from pathlib import Path
    import io

    ext = Path(filename).suffix.lower()
    result = {"text": "", "tables": [], "metadata": {}}

    try:
        if ext == ".docx":
            result = _read_docx(content, filename)
        elif ext == ".pptx":
            result = _read_pptx(content, filename)
        elif ext == ".pdf":
            result = _read_pdf(content, filename)
        elif ext in (".txt", ".md"):
            result = _read_text(content, filename)
        else:
            logger.warning(f"未対応のファイル形式: {ext}")
    except Exception as e:
        logger.exception(f"ドキュメント読み込みエラー: {e}")
        result["text"] = f"ドキュメント読み込みエラー: {e}"

    return result


def _read_docx(content: bytes | None, filename: str) -> dict:
    """Word (.docx) ファイルを読み込む。"""
    try:
        from docx import Document
    except ImportError:
        return {"text": "python-docx がインストールされていません。pip install python-docx を実行してください。",
                "tables": [], "metadata": {}}

    from pathlib import Path
    import io

    if content is not None:
        buf = io.BytesIO(content)
        doc = Document(buf)
    else:
        doc = Document(filename)

    text_parts = []
    tables = []

    # 段落のテキスト抽出
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    # テーブルの抽出
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if rows:
            df = pd.DataFrame(rows[1:], columns=rows[0] if len(rows) > 1 else None)
            tables.append(df)

    full_text = "\n".join(text_parts)

    return {
        "text": full_text,
        "tables": tables,
        "metadata": {
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables),
        }
    }


def _read_pptx(content: bytes | None, filename: str) -> dict:
    """PowerPoint (.pptx) ファイルを読み込む。"""
    try:
        from pptx import Presentation
    except ImportError:
        return {"text": "python-pptx がインストールされていません。pip install python-pptx を実行してください。",
                "tables": [], "metadata": {}}

    from pathlib import Path
    import io

    if content is not None:
        buf = io.BytesIO(content)
        prs = Presentation(buf)
    else:
        prs = Presentation(filename)

    text_parts = []
    tables = []

    for slide_num, slide in enumerate(prs.slides, 1):
        text_parts.append(f"--- スライド {slide_num} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    df = pd.DataFrame(rows[1:], columns=rows[0] if len(rows) > 1 else None)
                    tables.append(df)

    full_text = "\n".join(text_parts)

    return {
        "text": full_text,
        "tables": tables,
        "metadata": {
            "slides": len(prs.slides),
            "tables": len(tables),
        }
    }


def _read_pdf(content: bytes | None, filename: str) -> dict:
    """PDF ファイルからテキストを抽出する。"""
    try:
        import pdfplumber
    except ImportError:
        return {"text": "pdfplumber がインストールされていません。pip install pdfplumber を実行してください。",
                "tables": [], "metadata": {}}

    from pathlib import Path
    import io

    text_parts = []
    tables = []
    metadata = {"pages": 0}

    if content is not None:
        buf = io.BytesIO(content)
        with pdfplumber.open(buf) as pdf:
            metadata["pages"] = len(pdf.pages)
            for page in pdf.pages:
                if page.text:
                    text_parts.append(page.text)
                # テーブル抽出
                page_tables = page.extract_tables()
                for t in page_tables:
                    if t:
                        df = pd.DataFrame(t[1:], columns=t[0] if len(t) > 1 else None)
                        tables.append(df)
    else:
        with pdfplumber.open(filename) as pdf:
            metadata["pages"] = len(pdf.pages)
            for page in pdf.pages:
                if page.text:
                    text_parts.append(page.text)
                page_tables = page.extract_tables()
                for t in page_tables:
                    if t:
                        df = pd.DataFrame(t[1:], columns=t[0] if len(t) > 1 else None)
                        tables.append(df)

    full_text = "\n".join(text_parts)

    return {
        "text": full_text,
        "tables": tables,
        "metadata": metadata
    }


def _read_text(content: bytes | None, filename: str) -> dict:
    """テキストファイル (.txt, .md) を読み込む。"""
    from pathlib import Path

    if content is not None:
        text = content.decode("utf-8", errors="ignore")
    else:
        text = Path(filename).read_text(encoding="utf-8", errors="ignore")

    return {
        "text": text,
        "tables": [],
        "metadata": {"lines": len(text.splitlines())}
    }


class LLMDataAnalyst:
    """
    データ分析対話を管理するサービスクラス。
    GGUFProviderを使用してLLM推論を実行する。
    """

    def __init__(self):
        self.conversation_history: list[dict] = []
        self.last_suggestions: dict = {}
        self._data_summary: str = ""
        self._state_snapshot: dict = {}

    def build_data_summary(self, state: dict) -> str:
        """
        データの要約テキストを生成する。
        行数・列数・データ型・欠損・目的変数統計・相関係数・カージナリティなど。
        """
        df = state.get("df")
        if df is None:
            return "データが読み込まれていません。"

        lines = []
        lines.append(f"## データ概要")
        lines.append(f"- 行数: {len(df):,}行")
        lines.append(f"- 列数: {len(df.columns)}列")
        # データ種別の推定
        smiles_col = state.get("smiles_col", "")
        if smiles_col and smiles_col in df.columns:
            lines.append("- データ種別: SMILES（化学）データ")
        else:
            lines.append("- データ種別: 表（タブル）データ")
        lines.append("")

        # 列情報（カージナリティ分析含む）
        lines.append("## 列一覧・カージナリティ分析")
        target_col = state.get("target_col", "")
        exclude_cols = set(state.get("exclude_cols", []))
        numeric_cols = []
        categorical_cols = []

        for col in df.columns:
            dtype = str(df[col].dtype)
            na_count = int(df[col].isna().sum())
            na_pct = na_count / len(df) * 100 if len(df) > 0 else 0
            n_unique = int(df[col].nunique(dropna=True))

            role = ""
            if col == target_col:
                role = " [目的変数]"
            elif col == smiles_col:
                role = " [SMILES列]"
                role += " （記述子計算の候補）"
            elif col in exclude_cols:
                role = " [除外列]"

            # カージナリティの評価
            cardinality = ""
            if n_unique == 0:
                cardinality = "（全て欠損）"
            elif n_unique == 1:
                cardinality = "（定数値）"
            elif n_unique < 10:
                cardinality = f"（低カージナリティ: カテゴリ変数の可能性）"
                categorical_cols.append(col)
            elif n_unique > len(df) * 0.1:
                cardinality = f"（高カージナリティ: 一意値数が行数の10%超過）"
                if pd.api.types.is_numeric_dtype(df[col]):
                    numeric_cols.append(col)
            else:
                if pd.api.types.is_numeric_dtype(df[col]):
                    numeric_cols.append(col)
                else:
                    categorical_cols.append(col)

            lines.append(f"- {col}{role}: {dtype}, 欠損={na_count}件({na_pct:.1f}%), 一意値数={n_unique} {cardinality}")

        lines.append("")

        # カージナリティまとめ
        lines.append("## カージナリティまとめ")
        lines.append(f"- 数値列: {len(numeric_cols)}個 -> {numeric_cols}")
        lines.append(f"- カテゴリ列: {len(categorical_cols)}個 -> {categorical_cols}")
        lines.append("")

        # 目的変数の統計
        if target_col and target_col in df.columns:
            lines.append("## 目的変数の統計")
            tc = df[target_col]
            if pd and hasattr(pd, "api") and pd.api.types.is_numeric_dtype(tc):
                lines.append(f"- 平均: {tc.mean():.4f}")
                lines.append(f"- 中央値: {tc.median():.4f}")
                lines.append(f"- 標準偏差: {tc.std():.4f}")
                lines.append(f"- 最小値: {tc.min():.4f}")
                lines.append(f"- 最大値: {tc.max():.4f}")
                # 歪度・尖度
                try:
                    lines.append(f"- 歪度(skewness): {tc.skew():.4f}")
                    lines.append(f"- 尖度(kurtosis): {tc.kurtosis():.4f}")
                except Exception:
                    pass
            else:
                top_vals = tc.value_counts().head(5)
                lines.append(f"- クラス数: {tc.nunique()}")
                lines.append(f"- 最頻値: {dict(top_vals.iloc[:3])}")
                # クラス不均衡チェック
                if tc.nunique() <= 10:
                    val_counts = tc.value_counts()
                    min_pct = val_counts.min() / val_counts.sum() * 100
                    lines.append(f"- 最小クラス割合: {min_pct:.1f}% {'⚠️ 不均衡データ' if min_pct < 10 else ''}")
            lines.append("")

        # 相関係数分析（数値列のみ）
        if len(numeric_cols) >= 2 and target_col in numeric_cols:
            lines.append("## 相関係数分析（目的変数との相関）")
            try:
                corr_series = df[numeric_cols].corr()[target_col].drop(target_col)
                # 強相関・中相関・弱相関に分類
                strong = [(c, v) for c, v in corr_series.items() if abs(v) >= 0.7]
                moderate = [(c, v) for c, v in corr_series.items() if 0.3 <= abs(v) < 0.7]
                weak = [(c, v) for c, v in corr_series.items() if abs(v) < 0.3]

                if strong:
                    lines.append("- 強相関(|r|>=0.7):")
                    for c, v in sorted(strong, key=lambda x: abs(x[1]), reverse=True)[:5]:
                        lines.append(f"  * {c}: {v:.4f}")
                if moderate:
                    lines.append("- 中相関(0.3<=|r|<0.7):")
                    for c, v in sorted(moderate, key=lambda x: abs(x[1]), reverse=True)[:5]:
                        lines.append(f"  * {c}: {v:.4f}")
                if weak:
                    lines.append(f"- 弱相関(|r|<0.3): {len(weak)}個の変数")
                # 多重共線性チェック
                corr_matrix = df[numeric_cols].corr()
                high_corr_pairs = []
                for i, c1 in enumerate(numeric_cols):
                    for c2 in numeric_cols[i+1:]:
                        if abs(corr_matrix.loc[c1, c2]) >= 0.8:
                            high_corr_pairs.append((c1, c2, corr_matrix.loc[c1, c2]))
                if high_corr_pairs:
                    lines.append("- ⚠️ 多重共線性の可能性がある変数ペア:")
                    for c1, c2, v in high_corr_pairs[:5]:
                        lines.append(f"  * {c1} vs {c2}: {v:.4f}")
            except Exception as e:
                lines.append(f"- 相関係数計算エラー: {e}")
            lines.append("")

        # SMILES列の情報・分子特性分析
        if smiles_col and smiles_col in df.columns:
            lines.append("## SMILES列の情報・分子特性")
            smiles_list = df[smiles_col].dropna()
            lines.append(f"- SMILES件数: {len(smiles_list)}件")
            # 簡易バリデーション
            try:
                from rdkit import Chem
                valid = sum(1 for s in smiles_list if Chem.MolFromSmiles(str(s)) is not None)
                lines.append(f"- 有効SMILES: {valid}/{len(smiles_list)}件")
                # 分子特性の簡易計算
                if valid > 0:
                    mw_list, logp_list, hba_list, hbd_list = [], [], [], []
                    ring_list = []
                    for s in list(smiles_list)[:min(100, len(smiles_list))]:
                        mol = Chem.MolFromSmiles(str(s))
                        if mol:
                            try:
                                from rdkit.Chem import Descriptors
                                mw_list.append(Descriptors.MolWt(mol))
                                logp_list.append(Descriptors.MolLogP(mol))
                                hba_list.append(Descriptors.NumHAcceptors(mol))
                                hbd_list.append(Descriptors.NumHDonors(mol))
                                ring_list.append(mol.GetRingInfo().NumRings())
                            except Exception:
                                pass
                    if mw_list:
                        lines.append(f"- 分子量(MW): 平均={sum(mw_list)/len(mw_list):.1f}, 範囲={min(mw_list):.1f}~{max(mw_list):.1f}")
                    if logp_list:
                        lines.append(f"- LogP: 平均={sum(logp_list)/len(logp_list):.2f}, 範囲={min(logp_list):.2f}~{max(logp_list):.2f}")
                    if hba_list:
                        lines.append(f"- H受容体数: 平均={sum(hba_list)/len(hba_list):.1f}")
                    if hbd_list:
                        lines.append(f"- H供与体数: 平均={sum(hbd_list)/len(hbd_list):.1f}")
                    if ring_list:
                        lines.append(f"- 環数: 平均={sum(ring_list)/len(ring_list):.1f}")
                    # 推奨記述子
                    lines.append("- 推奨記述子: RDKit基本記述子を推奨（一般化学）")
                    if target_col in df.columns and pd.api.types.is_numeric_dtype(df[target_col]):
                        lines.append("  数値目的変数の場合: ECFP、物理化学記述子を推奨")
                    else:
                        lines.append("  分類目的変数の場合: ECFP、MACCS、PubChem記述子を推奨")
            except Exception:
                pass
            lines.append("")

        # タスクタイプ
        task_type = state.get("task_type", "auto")
        lines.append(f"## 推定タスク")
        lines.append(f"- タスクタイプ: {task_type}")
        if len(df) < 500:
            lines.append("- データサイズ: 少量（<500サンプル）→ 線形モデル・SVR等を推奨")
        elif len(df) < 5000:
            lines.append("- データサイズ: 中量（500-5000）→ RandomForest・XGBoost等を推奨")
        else:
            lines.append("- データサイズ: 大量（>5000）→ Neural Network・Stacking等を推奨")
        lines.append("")

        # データプレビュー（最初の3行）
        lines.append("## データプレビュー（最初の3行）")
        try:
            preview = df.head(3).to_string(index=False, max_colwidth=20)
            lines.append(f"```\n{preview}\n```")
        except Exception:
            pass

        summary = "\n".join(lines)
        self._data_summary = summary
        return summary

    def analyze(
        self,
        state: dict,
        user_message: str | None = None,
        reset: bool = False,
    ) -> dict:
        """
        データを分析し、LLMの返信を取得する。

        Args:
            state: アプリケーション状態
            user_message: ユーザーからの追加メッセージ（なければ初期分析）
            reset: True の場合、対話履歴をリセットする

        Returns:
            辞書: {"reply": str, "conversation": list, "suggestions": dict}
        """
        if reset:
            self.conversation_history = []

        # データ要約を生成（初回またはデータ変更時）
        current_summary = self.build_data_summary(state)
        self._state_snapshot = self._snapshot_state(state)

        # システムプロンプト
        system_prompt = self._build_system_prompt()

        # ユーザーメッセージの構築
        # ドキュメント内容の取得
        document_text = state.get("document_text", "")
        document_meta = state.get("document_metadata", {})

        if not self.conversation_history:
            # 初回：データ要約を含める
            parts = ["以下のデータについて分析プランを提案してください。"]
            if document_text:
                parts.append(f"\n\n## ドキュメント内容\n{document_text[:2000]}")  # 最初の2000文字
                if document_meta:
                    parts.append(f"\n\nドキュメントメタデータ: {document_meta}")
            if user_message:
                parts.append(f"\n\n{user_message}")
            parts.append(f"\n\n{current_summary}")
            user_content = "\n".join(parts)
        else:
            # 2回目以降：ユーザーメッセージのみ
            user_content = user_message or "補充情報はありません。分析を続けてください。"

        # 対話履歴に追加
        self.conversation_history.append({"role": "user", "content": user_content})

        # LLM呼び出し
        try:
            reply = self._call_llm(system_prompt, self.conversation_history)
            self.conversation_history.append({"role": "assitant", "content": reply})

            # 提案を解析
            suggestions = self._parse_suggestions(reply)

            return {
                "reply": reply,
                "conversation": list(self.conversation_history),
                "suggestions": suggestions,
                "data_summary": self._data_summary,
            }
        except Exception as e:
            logger.exception(f"[LLMDataAnalyst] LLM呼び出しエラー: {e}")
            error_reply = f"LLMの呼び出しに失敗しました: {e}"
            return {
                "reply": error_reply,
                "conversation": list(self.conversation_history),
                "suggestions": {},
                "data_summary": self._data_summary,
                "error": str(e),
            }

    def apply_suggestions(self, state: dict, suggestions: dict | None = None) -> list[str]:
        """
        LLMの提案をstateに反映する。
        反映した項目のリストを返す。
        """
        if suggestions is None:
            suggestions = self.last_suggestions

        applied = []
        logger.info(f"[LLMDataAnalyst] Applying suggestions: {suggestions}")

        # タスクタイプ
        if "task_type" in suggestions:
            val = suggestions["task_type"]
            if val in ("regression", "classification"):
                state["task_type"] = val
                applied.append(f"タスクタイプ: {val}")

        # 目的変数
        if "target_col" in suggestions:
            col = suggestions["target_col"]
            if col in state.get("df", {}).columns:
                state["target_col"] = col
                applied.append(f"目的変数: {col}")

        # 除外列
        if "exclude_cols" in suggestions:
            cols = suggestions["exclude_cols"]
            if isinstance(cols, list):
                state["exclude_cols"] = [c for c in cols if c in state.get("df", {}).columns]
                applied.append(f"除外列: {state['exclude_cols']}")

        # 選択モデル
        if "selected_models" in suggestions:
            models = suggestions["selected_models"]
            if isinstance(models, list):
                state["selected_models"] = models
                applied.append(f"選択モデル: {len(models)}件")

        # 数値スケーラー
        if "num_scaler" in suggestions:
            val = suggestions["num_scaler"]
            if val in ("standard", "robust", "minmax", "maxabs", "none"):
                state["num_scaler"] = val
                applied.append(f"数値スケーラー: {val}")

        # 特徴量選択
        if "feature_selector" in suggestions:
            val = suggestions["feature_selector"]
            if val in (
                "none", "variance", "selectkbest_f", "selectkbest_mi",
                "select_from_model_lasso", "select_from_model_rf", "rfe", "boruta",
            ):
                state["feature_selector"] = val
                applied.append(f"特徴量選択: {val}")

        # SMILES記述子
        if "selected_descriptors" in suggestions:
            descs = suggestions["selected_descriptors"]
            if isinstance(descs, list):
                state["selected_descriptors"] = descs
                applied.append(f"記述子: {len(descs)}件")

        # 単調性制約
        if "monotonic_constraints" in suggestions:
            constraints = suggestions["monotonic_constraints"]
            if isinstance(constraints, dict):
                # state["monotonic_constraints"] に反映 (data_tab.py 形式)
                state["monotonic_constraints"] = {}
                for col, direction in constraints.items():
                    if direction == 1:
                        state["monotonic_constraints"][col] = 1
                    elif direction == -1:
                        state["monotonic_constraints"][col] = -1
                    # direction == 0 は制約なし

                # state["feature_constraints"] にも反映 (monotonicity_config.py 形式)
                if "feature_constraints" not in state:
                    state["feature_constraints"] = {}
                for col, direction in constraints.items():
                    if direction == 1:
                        state["feature_constraints"][col] = {
                            "direction": "increasing",
                            "strength": 1.0,
                            "sigma": 3.0
                        }
                    elif direction == -1:
                        state["feature_constraints"][col] = {
                            "direction": "decreasing",
                            "strength": 1.0,
                            "sigma": 3.0
                        }
                    elif col in state["feature_constraints"]:
                        del state["feature_constraints"][col]

                applied.append(f"単調性制約: {len([c for c in constraints.values() if c != 0])}件")

        self.last_suggestions = suggestions
        return applied

    def _build_system_prompt(self) -> str:
        """システムプロンプトを構築する。"""
        return """あなたは化学データ分析の専門家です。ユーザーがアップロードしたデータを分析し、最適な機械学習プランを一緒に立てます。

重要な指示:
1. データ要約をよく読み、以下の点を確認してください:
   - 各列の意味や役割（SMILES列、目的変数、特徴量など）
   - データの品質（欠損値、一意値数、データ型）
   - タスクタイプ（回帰か分類か）
   - サンプル数が少ない場合の交差検証戦略

2. 以下の場合は必ずユーザーに質問してください:
   - 列の意味が不明確な場合（例: 「この'Property'列は何を表していますか？」）
   - 分析目的が明確でない場合（例: 「このデータで予測したい対象は何ですか？」）
   - 目的変数の選択で複数の候補がある場合
   - データの品質に問題がある場合（例: 「欠損値が30%以上ある列がありますが、どう処理しますか？」）
   - SMILES列の有無や、記述子計算が必要かどうか不明な場合

3. ドキュメント分析（Word/PowerPoint/PDF等）:
   - ドキュメント内容が提供された場合は、その内容をよく読み込んでください
   - ドキュメントから以下の情報を抽出してください:
     * 分析目的（予測したい物性は何か？）
     * データの背景（実験条件、サンプル情報など）
     * 推奨される特徴量や記述子
     * 既知の相関関係や制約条件
   - ドキュメントの情報を基に、単調性制約や特徴量選択の推奨を行ってください
   - 表データの場合: ドキュメントから各変数の物理化学的意味を理解し、単調性制約を推奨
   - SMILESデータの場合: ドキュメントから予測対象の物性を理解し、適切な記述子を推奨

4. 推奨モデルはデータサイズに応じて選択してください:
   - 少量データ（<500サンプル）: シンプルなモデル（Linear、Ridge、SVRなど）
   - 中量データ（500-5000）: 標準的なモデル（RandomForest、XGBoostなど）
   - 大量データ（>5000）: 高度なモデル（Neural Network、Stackingなど）

5. 提案するプランは具体的に:
   - タスクタイプ（回帰/分類）とその理由
   - 目的変数の選択理由
   - 推奨モデルとその理由
   - 特徴量エンジニアリングの提案
   - SMILES記述子が必要かどうか
   - 単調性制約の推奨（表データの場合、ドキュメント情報も考慮）

6. 回答形式:
   - 最初にデータの概要を簡潔に説明
   - ドキュメントがある場合は、その要約と分析への活用方法を説明
   - 不明点があれば質問（質問は明確で具体的に）
   - 最後に以下のJSON形式で提案を出力（コードブロック内）:

```json
{
  "task_type": "regression" または "classification",
  "target_col": "目的変数名",
  "exclude_cols": ["除外列1", "除外列2"],
  "selected_models": ["モデルキー1", "モデルキー2"],
  "num_scaler": "standard" など,
  "feature_selector": "none" など,
  "selected_descriptors": ["記述子1", "記述子2"],
  "monotonic_constraints": {"変数名": 1 または -1},  // 1=増加, -1=減少, 0=制約なし
  "notes": "その他のアドバイス"
}
```

- JSONは必ず出力すること。不明点がある場合は、JSONの前に質問文を書き、ユーザーの返答を待ってからJSONを出力すること。
- 日本語で回答することを忘れないでください。
"""

    def _call_llm(self, system_prompt: str, conversation_history: list[dict]) -> str:
        """LLMを呼び出す。OpenAI互換APIを優先。"""
        # Try OpenAI-compatible API first (Ollama etc.)
        try:
            from backend.llm.providers.openai_provider import OpenAIProvider, load_openai_config
            cfg = load_openai_config()
            if cfg.get("base_url"):
                provider = OpenAIProvider()
                return self._call_openai(provider, system_prompt, conversation_history)
        except Exception as e:
            logger.info(f"[LLMDataAnalyst] OpenAIProvider unavailable: {e}")

        # Fallback: GGUFProvider
        try:
            from backend.llm.providers.gguf_provider import (
                GGUFProvider, load_gguf_config, load_gguf_model,
            )
        except ImportError as e:
            raise RuntimeError(f"GGUFProvider not available: {e}")

        cfg = load_gguf_config()
        model_id = cfg.get("model_id", "prism-ml/Bonsai-8B-gguf")
        filename = cfg.get("filename", "Bonsai-8B-Q1_0.gguf")
        n_ctx = cfg.get("n_ctx", 4096)

        # Build prompt
        prompt_parts = []
        if system_prompt:
            prompt_parts.append("<|im_start|>system\n" + system_prompt + "\n<|im_end|>")
        max_chars = (min(n_ctx, 2048) - 512) * 4
        current_len = 0
        truncated_history = []
        for msg in reversed(conversation_history):
            role = msg.get("role", "user")
            content_str = msg.get("content", "")
            line = ""
            if role == "user":
                line = "<|im_start|>user\n" + content_str + "\n<|im_end|>"
            elif role == "assitant":
                line = "<|im_start|>assitant\n" + content_str + "\n<|im_end|>"
            if line:
                if current_len + len(line) > max_chars:
                    break
                truncated_history.append(line)
                current_len += len(line)
        truncated_history.reverse()
        prompt_parts.extend(truncated_history)
        prompt_parts.append("<|im_start|>assitant\n")
        full_prompt = "".join(prompt_parts)

        provider = GGUFProvider(model_id=model_id, filename=filename)
        model_path = provider._resolve_model_path()
        model = load_gguf_model(model_path, n_gpu_layers=0, n_ctx=min(n_ctx, 2048), n_batch=512)
        logger.info(f"[LLMDataAnalyst] Prompt length (chars): {len(full_prompt)}")
        result = model(full_prompt, max_tokens=2048, temperature=0.3, echo=False)
        return result["choices"][0]["text"].strip()

    def _call_openai(self, provider, system_prompt: str, conversation_history: list[dict]) -> str:
        """Call OpenAI-compatible API (Ollama etc.)"""
        messages = []
        if system_prompt:
            messages.append({"role": "system", "content": system_prompt})
        for msg in conversation_history:
            role = msg.get("role", "user")
            content = msg.get("content", "")
            if role in ("user", "assitant"):
                messages.append({"role": role, "content": content})
        messages.append({"role": "assitant", "content": ""})

        logger.info(f"[LLMDataAnalyst] OpenAI API call with {len(messages)} messages")
        request = type("LLMRequest", (), {
            "system_prompt": system_prompt,
            "user_prompt": conversation_history[-1].get("content", "") if conversation_history else "",
            "max_tokens": 2048,
            "temperature": 0.3,
        })()
        response = provider.generate(request)
        return response.content.strip()

    def _snapshot_state(self, state: dict) -> dict:
        """状態のスナップショットを作成する。"""
        import copy
        return copy.deepcopy(state)

    def _parse_suggestions(self, reply: str) -> dict:
        """LLMの返信から提案を解析する。"""
        import re
        import json

        suggestions = {}

        # JSONブロックを探す
        json_match = re.search(r"```json\s*(.*?)\s*```", reply, re.DOTALL)
        if json_match:
            try:
                data = json.loads(json_match.group(1))
                suggestions.update(data)
            except json.JSONDecodeError:
                pass

        return suggestions

    def recommend_descriptors(
        self,
        state: dict,
        target_description: str = "",
        interview_notes: str = "",
    ) -> dict:
        """
        目的変数に対して化学物性に基づく記述子を推奨する。

        Args:
            state: アプリケーション状態（df, target_col, smiles_col等を含む）
            target_description: ユーザーからの目的変数の追加説明
            interview_notes: ヒヤリングメモ

        Returns:
            辞書: { "reply": str, "recommendations": dict, "conversation": list }
        """
        df = state.get("df")
        target_col = state.get("target_col", "")
        smiles_col = state.get("smiles_col", "")

        if df is None or not target_col:
            return {
                "reply": "データまたは目的変数が設定されていません。",
                "recommendations": {},
                "conversation": list(self.conversation_history),
                "error": "missing_data",
            }

        # データサンプルを作成
        sample_df = df.head(5)
        df_sample = sample_df.to_string(index=False, max_colwidth=30)

        # 列リスト
        df_columns = list(df.columns)

        # プロンプトを構築
        prompt = build_descriptor_recommendation_prompt(
            target_col=target_col,
            target_description=target_description,
            df_columns=df_columns,
            df_sample=df_sample,
            interview_notes=interview_notes,
        )

        # システムプロンプト
        system_prompt = """あなたは化学物性予測のための記述子選択の専門家です。
化学的知見に基づいて、目的変数を予測するための最適な記述子を推奨してください。

重要な指示:
1. ユーザーが予測したい物性を理解し、その物性に最も寄与する分子特性を考慮する
2. 既存の列で同じような記述子がある場合は、それらを優先的に使用する
3. SMILES列がある場合は、RDKit記述子やフィンガープリントを推奨する
4. 各記述子について、なぜその物性に寄与するかを化学的根拠と共に説明する
5. 最大10個までの記述子を推奨する（優先度順）
6. 物性が不明確な場合は、ユーザーに質問を投げる
7. 日本語で回答する
"""

        # 対話履歴に追加
        self.conversation_history.append({"role": "user", "content": prompt})

        try:
            reply = self._call_llm(system_prompt, self.conversation_history)
            self.conversation_history.append({"role": "assitant", "content": reply})

            # 推奨結果を解析
            recommendations = self._parse_descriptor_recommendations(reply)

            return {
                "reply": reply,
                "recommendations": recommendations,
                "conversation": list(self.conversation_history),
                "matched_properties": find_matching_properties(target_col, target_description)[:3],
            }
        except Exception as e:
            logger.exception(f"[LLMDataAnalyst] 記述子推奨エラー: {e}")
            return {
                "reply": f"記述子推奨でエラーが発生しました: {e}",
                "recommendations": {},
                "conversation": list(self.conversation_history),
                "error": str(e),
            }

    def _parse_descriptor_recommendations(self, reply: str) -> dict:
        """
        LLMの返信から記述子推奨情報を解析する。
        JSONブロックまたは構造化テキストから情報を抽出。
        """
        import re
        import json

        result = {
            "descriptors": [],
            "questions": [],
            "matched_property": "",
            "confidence": "low",
            "notes": "",
        }

        # JSONブロックを探す
        json_match = re.search(r"```json\s*(.*?)\s*```", reply, re.DOTALL)
        if json_match:
            try:
                data = json.loads(json_match.group(1))
                result["descriptors"] = data.get("recommended_descriptors", [])
                result["questions"] = data.get("interview_questions", [])
                result["matched_property"] = data.get("matched_property", "")
                result["confidence"] = data.get("confidence", "low")
                result["notes"] = data.get("notes", "")
                return result
            except json.JSONDecodeError:
                pass

        # JSONがない場合はテキストから抽出を試みる
        lines = reply.split("\n")
        for line in lines:
            line = line.strip()
            # 記述子名を探す（「- 」や「* 」で始まる行）
            if line.startswith(("- ", "* ", "• ")):
                # 記述子のパターン: 名前 (ソース): 理由
                match = re.match(r"[-*•]\s*(\w+)\s*\((\w+)\):\s*(.*)", line)
                if match:
                    result["descriptors"].append({
                        "name": match.group(1),
                        "source": match.group(2),
                        "reason": match.group(3),
                        "priority": len(result["descriptors"]) + 1,
                    })
            # 質問を探す
            if "？" in line or "?" in line:
                if line not in result["questions"]:
                    result["questions"].append(line)

        return result

    def start_interview(
        self,
        state: dict,
        user_response: str = "",
        interview_history: list[dict] | None = None,
    ) -> dict:
        """
        ユーザーとのヒヤリング（対話）を通じて記述子を推奨する。

        Args:
            state: アプリケーション状態
            user_response: ユーザーからの返答（空なら初回）
            interview_history: これまでのヒヤリング履歴

        Returns:
            辞書: { "reply": str, "interview_complete": bool, "recommendations": dict }
        """
        if interview_history is None:
            interview_history = []

        system_prompt = """あなたは化学データ分析の専門家です。
ユーザーが予測したい物性についてヒヤリングを行い、最適な記述子を推奨してください。

ヒヤリングの手順:
1. 目的変数の物性が明確でない場合は、積極的に質問して情報を引き出してください:
   - 「この目的変数はどのような物性を表していますか？」（例: 屈折率、溶解度、融点、沸点、極性など）
   - 「この物性に影響すると考えられる分子の特性は何ですか？」（例: 分子量、極性、官能基、芳香環の有無など）
   - 「Similarな化合物でこの物性が高い/低い例はありますか？これによりどのような構造的特徴が重要か推測できますか？」
   - 「この物性を予測する上で、どの程度の精度が必要ですか？」（スクリーニング用か、定量予測か）

2. 初回のヒヤリングでは、以下の基本情報も合わせて確認してください:
   - データセットのサイズと特徴（サンプル数は十分か？）
   - 目的変数のデータタイプ（連続値か離散値か）
   - 欠損値や外れ値の状況

3. 物性が比較的特定できたら、以下の形式で記述子を推奨してください:

```json
{
  "interview_complete": true,
  "matched_property": "特定された物性名",
  "recommended_descriptors": [
    {"name": "記述子名", "source": "RDKit", "reason": "推奨理由", "priority": 1}
  ],
  "notes": "アドバイス"
}
```

4. まだ情報が不十分な場合は、interview_complete: false として、具体的な質問を続けてください。
5. 日本語で回答してください。
6. 質問は1回につき2-3個に絞り、ユーザーの負担を減らしてください。
7. 前の回答を踏まえて、次の質問につながるように会話を導いてください。
"""

        # プロンプト構築
        df = state.get("df")
        target_col = state.get("target_col", "")

        if not user_response:
            # 初回：ターゲット情報を提示して質問開始
            prompt_parts = [
                f"目的変数: {target_col}",
                f"データ列: {', '.join(df.columns) if df is not None else 'なし'}",
                "",
                "この目的変数について、以下の点を教えていただけますか？",
                "1. この目的変数はどのような物性を表していますか？（例: 屈折率、溶解度、融点、沸点、極性、反応速度など）",
                "2. この物性に影響すると考えられる分子の特性は何ですか？（例: 分子量、極性、官能基、芳香環の有無、立体構造など）",
                "3. この物性を予測する目的は何ですか？（例: 新規化合物のスクリーニング、構造最適化、メカニズム解明など）",
                "",
                "可能な範囲でお答えください。不明点は「わからない」と答えていただければ、こちらで推測して質問を続けます。",
            ]
            prompt = "\n".join(prompt_parts)
        else:
            # 2回目以降：ユーザーの返答を受けて継続
            prompt = f"ユーザーの返答: {user_response}\n\n次の質問をするか、記述子を推奨してください。"

        interview_history.append({"role": "user", "content": prompt})

        try:
            reply = self._call_llm(system_prompt, interview_history)
            interview_history.append({"role": "assitant", "content": reply})

            # 結果を解析
            import re
            import json

            result = {
                "reply": reply,
                "interview_history": list(interview_history),
                "interview_complete": False,
                "recommendations": {},
            }

            # JSONを探す
            json_match = re.search(r"```json\s*(.*?)\s*```", reply, re.DOTALL)
            if json_match:
                try:
                    data = json.loads(json_match.group(1))
                    result["interview_complete"] = data.get("interview_complete", False)
                    result["recommendations"] = {
                        "descriptors": data.get("recommended_descriptors", []),
                        "matched_property": data.get("matched_property", ""),
                        "notes": data.get("notes", ""),
                    }
                except json.JSONDecodeError:
                    pass

            return result
        except Exception as e:
            logger.exception(f"[LLMDataAnalyst] ヒヤリングエラー: {e}")
            return {
                "reply": f"ヒヤリングでエラー: {e}",
                "interview_history": list(interview_history),
                "interview_complete": False,
                "error": str(e),
            }

    def reset(self) -> None:
        """対話履歴と状態をリセットする。"""
        self.conversation_history = []
        self.last_suggestions = {}
        self._data_summary = ""
        self._state_snapshot = {}


# Global instance
_global_analyst = None


def get_data_analyst() -> LLMDataAnalyst:
    """グローバルLLMDataAnalystインスタンスを取得する。"""
    global _global_analyst
    if _global_analyst is None:
        _global_analyst = LLMDataAnalyst()
    return _global_analyst


def reset_data_analyst() -> None:
    """グローバルインスタンスをリセットする。"""
    global _global_analyst
    if _global_analyst:
        _global_analyst.reset()
    _global_analyst = LLMDataAnalyst()
